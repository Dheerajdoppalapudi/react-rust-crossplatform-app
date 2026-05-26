"""
Export router — converts HTML slide presentations to PPTX.

Uses Playwright to screenshot each slide (navigating via #btn-next),
then assembles the screenshots into a python-pptx file (16:9 widescreen).

Fixes applied:
  L-6  : Rate-limited to 5 req/min per IP (Playwright is expensive).
  H-4  : _playwright_sem(2) prevents more than 2 concurrent browser instances
          per worker process, capping thread-pool exhaustion.
  H-4  : browser.close() is guaranteed via try/finally inside the playwright
          context manager so leaked browser processes can't accumulate.
"""

import asyncio
import structlog
import os
import tempfile
from io import BytesIO

from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from core.limiter import limiter

logger = structlog.get_logger(__name__)

router = APIRouter()

# At most 2 concurrent Playwright browser instances per worker.
# Each Chromium instance consumes ~100–200 MB RAM; 2 is safe on a 1 GB container.
# Raise via PPTX_CONCURRENCY env var once the worker fleet is sized for it.
import os as _os
_MAX_CONCURRENT = int(_os.getenv("PPTX_CONCURRENCY", "2"))
_playwright_sem = asyncio.Semaphore(_MAX_CONCURRENT)


class PptxExportRequest(BaseModel):
    html: str
    title: str = "Presentation"


@router.post("/export/pptx")
@limiter.limit("5/minute")
async def export_to_pptx(payload: PptxExportRequest, request: Request):
    """
    Convert an HTML slide presentation to a PPTX file.

    Playwright opens the HTML, screenshots each slide by clicking #btn-next,
    then python-pptx assembles the screenshots into a 16:9 widescreen deck.
    """
    async with _playwright_sem:
        try:
            pptx_bytes = await _render_pptx(payload.html, payload.title)
        except Exception as exc:
            logger.error("pptx_export_failed", error=str(exc), exc_info=True)
            raise HTTPException(status_code=500, detail="PPTX generation failed. Please try again.")

    safe = "".join(c for c in payload.title if c.isalnum() or c in " _-")[:50].strip()
    filename = f"{safe.replace(' ', '_') or 'presentation'}.pptx"

    return StreamingResponse(
        BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


async def _render_pptx(html: str, title: str) -> bytes:
    from playwright.async_api import async_playwright
    from pptx import Presentation
    from pptx.util import Inches

    with tempfile.NamedTemporaryFile(mode="w", suffix=".html", delete=False, encoding="utf-8") as f:
        f.write(html)
        tmp_path = f.name

    screenshots: list[bytes] = []
    try:
        async with async_playwright() as pw:
            browser = await pw.chromium.launch()
            try:
                page = await browser.new_page(viewport={"width": 1280, "height": 720})
                await page.goto(f"file://{tmp_path}")
                await page.wait_for_timeout(600)

                slide_count = await page.evaluate("""
                    () => {
                        const deck = document.getElementById('deck');
                        if (deck) {
                            const n = parseInt(deck.getAttribute('data-count'), 10);
                            if (n > 0) return n;
                        }
                        return document.querySelectorAll('.slide').length || 1;
                    }
                """)
                slide_count = max(1, min(int(slide_count), 50))

                for i in range(slide_count):
                    img = await page.screenshot(type="png", full_page=False)
                    screenshots.append(img)
                    if i < slide_count - 1:
                        try:
                            await page.click("#btn-next")
                            await page.wait_for_timeout(450)
                        except Exception:
                            break
            finally:
                await browser.close()
    finally:
        os.unlink(tmp_path)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for screenshot in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            BytesIO(screenshot),
            Inches(0), Inches(0),
            Inches(13.33), Inches(7.5),
        )

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
