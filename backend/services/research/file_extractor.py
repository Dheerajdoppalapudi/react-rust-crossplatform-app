"""
Extract plain text from uploaded documents and detect URLs in user messages.
"""

import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

_URL_RE = re.compile(r'https?://[^\s<>"\']+')


def extract_urls_from_text(text: str) -> list[str]:
    """Return all HTTP/HTTPS URLs found in the text."""
    return _URL_RE.findall(text or "")


def extract_text(file_path: str) -> str:
    """
    Extract plain text from a file based on its extension.
    Returns empty string on failure rather than raising.
    """
    path = Path(file_path)
    if not path.exists():
        logger.warning("file_extractor: path does not exist: %s", file_path)
        return ""

    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        if ext == ".pptx":
            return _extract_pptx(file_path)
        if ext in (".txt", ".md", ".markdown"):
            return path.read_text(encoding="utf-8", errors="ignore")
        logger.warning("file_extractor: unsupported extension %s", ext)
        return ""
    except Exception as e:
        logger.error("file_extractor: failed to extract %s: %s", file_path, e)
        return ""


def _extract_pdf(file_path: str) -> str:
    import pdfplumber
    pages: list[str] = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs if run.text).strip()
                if line:
                    parts.append(line)
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
    return "\n\n".join(slides)
